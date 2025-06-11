"""
Ferramentas para manipulação de arquivos do Microsoft Office
"""

import os
import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
import json

try:
    import docx
    import openpyxl
    import pptx
except ImportError:
    logging.warning("Bibliotecas do Office não encontradas. Algumas funcionalidades estarão indisponíveis.")

from ...core.logs import setup_logging
from ..exceptions import ToolError, FileError

# Configurar logger
logger = setup_logging(__name__)

def criar_word(caminho_arquivo: str, texto: str,
    mcp_client: Optional[Any] = None
) -> str:
    """
    Cria um arquivo Word (.docx) com o texto fornecido.
    
    Args:
        caminho_arquivo: Caminho onde o arquivo será salvo
        texto: Conteúdo a ser escrito no documento
        
    Returns:
        Mensagem de confirmação
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando criar_word para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "criar_word",
                "parametros": {"caminho_arquivo": caminho_arquivo, "texto": texto}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            return resultado.get("resultado", "Erro: Resultado vazio do MCP Server.")

        # Lógica existente (fallback local)
        doc = docx.Document()
        doc.add_paragraph(texto)
        doc.save(caminho_arquivo)
        return f"Arquivo Word criado com sucesso: {caminho_arquivo}"
    except Exception as e:
        logger.error(f"Erro ao criar arquivo Word: {str(e)}")
        raise ToolError(f"Erro ao criar arquivo Word: {str(e)}")

def ler_word(caminho_arquivo: str,
    mcp_client: Optional[Any] = None
) -> str:
    """
    Lê o texto de um arquivo Word (.docx).
    
    Args:
        caminho_arquivo: Caminho do arquivo Word
        
    Returns:
        Conteúdo do documento
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando ler_word para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "ler_word",
                "parametros": {"caminho_arquivo": caminho_arquivo}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            return resultado.get("resultado", "Erro: Resultado vazio do MCP Server.")

        # Lógica existente (fallback local)
        doc = docx.Document(caminho_arquivo)
        texto = "\n".join([p.text for p in doc.paragraphs])
        return texto
    except Exception as e:
        logger.error(f"Erro ao ler arquivo Word: {str(e)}")
        raise ToolError(f"Erro ao ler arquivo Word: {str(e)}")

def criar_excel(caminho_arquivo: str, dados: List[List[Any]],
    mcp_client: Optional[Any] = None
) -> str:
    """
    Cria um arquivo Excel (.xlsx) com os dados fornecidos.
    
    Args:
        caminho_arquivo: Caminho onde o arquivo será salvo
        dados: Lista de listas contendo os dados a serem escritos
        
    Returns:
        Mensagem de confirmação
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando criar_excel para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "criar_excel",
                "parametros": {"caminho_arquivo": caminho_arquivo, "dados": dados}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            return resultado.get("resultado", "Erro: Resultado vazio do MCP Server.")

        # Lógica existente (fallback local)
        wb = openpyxl.Workbook()
        ws = wb.active
        for linha in dados:
            ws.append(linha)
        wb.save(caminho_arquivo)
        return f"Arquivo Excel criado com sucesso: {caminho_arquivo}"
    except Exception as e:
        logger.error(f"Erro ao criar arquivo Excel: {str(e)}")
        raise ToolError(f"Erro ao criar arquivo Excel: {str(e)}")

def ler_excel(caminho_arquivo: str,
    mcp_client: Optional[Any] = None
) -> List[List[Any]]:
    """
    Lê o conteúdo de um arquivo Excel (.xlsx).
    
    Args:
        caminho_arquivo: Caminho do arquivo Excel
        
    Returns:
        Lista de listas contendo os dados lidos
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando ler_excel para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "ler_excel",
                "parametros": {"caminho_arquivo": caminho_arquivo}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            # O resultado esperado para ler_excel é uma lista de listas
            return resultado.get("resultado", []) # Retorna lista vazia em caso de erro ou resultado vazio

        # Lógica existente (fallback local)
        wb = openpyxl.load_workbook(caminho_arquivo)
        ws = wb.active
        dados = [[cell.value for cell in row] for row in ws.iter_rows()]
        return dados
    except Exception as e:
        logger.error(f"Erro ao ler arquivo Excel: {str(e)}")
        raise ToolError(f"Erro ao ler arquivo Excel: {str(e)}")

def criar_ppt(caminho_arquivo: str, titulos: List[str], textos: List[str],
    mcp_client: Optional[Any] = None
) -> str:
    """
    Cria um arquivo PowerPoint (.pptx) com slides de títulos e textos.
    
    Args:
        caminho_arquivo: Caminho onde o arquivo será salvo
        titulos: Lista de títulos para os slides
        textos: Lista de textos para os slides
        
    Returns:
        Mensagem de confirmação
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando criar_ppt para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "criar_ppt",
                "parametros": {"caminho_arquivo": caminho_arquivo, "titulos": titulos, "textos": textos}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            return resultado.get("resultado", "Erro: Resultado vazio do MCP Server.")

        # Lógica existente (fallback local)
        prs = pptx.Presentation()
        for titulo, texto in zip(titulos, textos):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = titulo
            slide.placeholders[1].text = texto
        prs.save(caminho_arquivo)
        return f"Arquivo PowerPoint criado com sucesso: {caminho_arquivo}"
    except Exception as e:
        logger.error(f"Erro ao criar arquivo PowerPoint: {str(e)}")
        raise ToolError(f"Erro ao criar arquivo PowerPoint: {str(e)}")

def ler_ppt(caminho_arquivo: str,
    mcp_client: Optional[Any] = None
) -> List[Dict[str, str]]:
    """
    Lê os títulos e textos dos slides de um arquivo PowerPoint (.pptx).
    
    Args:
        caminho_arquivo: Caminho do arquivo PowerPoint
        
    Returns:
        Lista de dicionários contendo título e texto de cada slide
    """
    try:
        # Refatoração para usar MCP Client se disponível
        if mcp_client:
            logger.info(f"Delegando ler_ppt para MCP Server para caminho: {caminho_arquivo}")
            tarefa_execucao = {
                "tipo": "executar_ferramenta",
                "nome_ferramenta": "ler_ppt",
                "parametros": {"caminho_arquivo": caminho_arquivo}
            }
            resultado = mcp_client.distribuir_tarefa(
                tarefa=json.dumps(tarefa_execucao),
                agente_id="Agente Local"
            )
            # O resultado esperado para ler_ppt é uma lista de dicionários
            return resultado.get("resultado", []) # Retorna lista vazia em caso de erro ou resultado vazio

        # Lógica existente (fallback local)
        prs = pptx.Presentation(caminho_arquivo)
        slides = []
        for slide in prs.slides:
            titulo = slide.shapes.title.text if slide.shapes.title else ""
            texto = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    texto += shape.text + "\n"
            slides.append({"titulo": titulo, "texto": texto.strip()})
        return slides
    except Exception as e:
        logger.error(f"Erro ao ler arquivo PowerPoint: {str(e)}")
        raise ToolError(f"Erro ao ler arquivo PowerPoint: {str(e)}") 