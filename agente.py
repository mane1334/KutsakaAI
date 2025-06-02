"""
Agente IA - Assistente Virtual Inteligente
Versão 1.0.0
"""

import os
import json
import logging
import sys
import importlib
import pkg_resources
from datetime import datetime
from typing import List, Dict, Any, Optional, Union, TypeVar, Generic, Set, Tuple
from dataclasses import dataclass, field
from enum import Enum, auto
from pathlib import Path
from contextlib import contextmanager
from functools import wraps
import traceback
import asyncio
from concurrent.futures import ThreadPoolExecutor
import signal
import psutil
import platform
import socket
import ssl
import certifi
from typing_extensions import TypeGuard, Protocol, runtime_checkable
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_community.chat_models import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain.agents import create_tool_calling_agent, AgentExecutor, initialize_agent, AgentType
from langchain.memory import ConversationBufferMemory
from ferramentas import ferramentas
from agenteia.core import gerenciador_gatilhos, TipoGatilho
import agenteia.config as config
from langchain_core.callbacks import BaseCallbackHandler
import re
from langchain.agents import create_openai_functions_agent
from langchain.tools import BaseTool
from langchain_community.chat_models import ChatOllama
import time
from .agenteia.config import CONFIG, validar_configuracoes
from langchain_core.utils.utils import secret_from_env
from pydantic import Field, SecretStr

# Configuração de logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('agente.log'),
        logging.StreamHandler()
    ]
)

logger = logging.getLogger(__name__)

# Tipos personalizados para validação
T = TypeVar('T')
class Resultado(Generic[T]):
    def __init__(self, valor: Optional[T] = None, erro: Optional[str] = None):
        self.valor = valor
        self.erro = erro
    
    @property
    def sucesso(self) -> bool:
        return self.erro is None
    
    @property
    def falha(self) -> bool:
        return self.erro is not None

class TipoMensagem(Enum):
    USUARIO = "usuario"
    AGENTE = "agente"
    SISTEMA = "sistema"

@dataclass
class Mensagem:
    tipo: TipoMensagem
    conteudo: str
    timestamp: datetime
    
    @classmethod
    def criar(cls, tipo: TipoMensagem, conteudo: str) -> 'Mensagem':
        if not conteudo or not conteudo.strip():
            raise ValueError("Conteúdo da mensagem não pode ser vazio")
        return cls(tipo=tipo, conteudo=conteudo.strip(), timestamp=datetime.now())

class ValidacaoError(Exception):
    """Exceção para erros de validação."""
    pass

class AgenteError(Exception):
    """Exceção personalizada para erros do agente."""
    pass

def validar_caminho(caminho: Union[str, Path]) -> Resultado[Path]:
    """Valida um caminho de arquivo."""
    try:
        path = Path(caminho)
        if not path.exists():
            return Resultado(erro=f"Caminho não existe: {caminho}")
        return Resultado(valor=path)
    except Exception as e:
        return Resultado(erro=f"Erro ao validar caminho: {str(e)}")

def validar_modelo(modelo: str) -> Resultado[str]:
    """Valida o nome do modelo."""
    modelos_validos = ["qwen3-4b", "llama-3.2-1b-instruct"]
    if modelo not in modelos_validos:
        return Resultado(erro=f"Modelo inválido: {modelo}. Modelos válidos: {modelos_validos}")
    return Resultado(valor=modelo)

def get_func_name(func):
    """Obtém o nome original da função, mesmo se decorada ou for um objeto tool."""
    # Caso seja um objeto tool do LangChain
    if hasattr(func, "func"):
        return get_func_name(func.func)
    # Caso seja decorada (lru_cache, wraps, etc)
    while hasattr(func, "__wrapped__"):
        func = func.__wrapped__
    return getattr(func, "__name__", None)

class ContextoManager:
    def __init__(self, max_mensagens: int = 100):
        self.max_mensagens = max_mensagens
        self.historico = []

    def adicionar_mensagem(self, tipo: str, conteudo: str):
        self.historico.append({
            'tipo': tipo,
            'mensagem': conteudo,
            'timestamp': datetime.now().isoformat()
        })
        if len(self.historico) > self.max_mensagens:
            self.historico.pop(0)

    def obter_historico(self):
        return self.historico

    def limpar_historico(self):
        self.historico = []

class PensamentoCallback(BaseCallbackHandler):
    """Callback para capturar o pensamento do agente em tempo real."""
    
    def __init__(self):
        self.pensamentos = []
        self.resposta_final = None
        self.em_pensamento = True
        self.ultima_ferramenta = None
    
    def on_llm_start(self, *args, **kwargs):
        """Chamado quando o LLM começa a gerar."""
        self.em_pensamento = True
        self.pensamentos = []
        self.ultima_ferramenta = None
    
    def on_llm_new_token(self, token: str, *args, **kwargs):
        """Chamado quando um novo token é gerado."""
        if token.strip():
            if self.em_pensamento:
                # Remove tags e marcadores indesejados
                token = token.replace("<think>", "").replace("</think>", "")
                token = re.sub(r">.*?<", "", token)  # Remove conteúdo entre > e <
                
                # Ignora tokens que são parte da resposta final ou pensamentos confusos
                if not any(token.startswith(m) for m in ["💭", "**Resposta:**", ">", "<", "Wait,", "Hmm,", "Okay,"]):
                    # Limpa o token de espaços extras e quebras de linha
                    token = token.strip()
                    if token and not token.isspace():
                        self.pensamentos.append(token)
    
    def on_chain_start(self, *args, **kwargs):
        """Chamado quando uma chain começa."""
        self.em_pensamento = True
    
    def on_chain_end(self, outputs, *args, **kwargs):
        """Chamado quando uma chain termina."""
        if isinstance(outputs, dict) and 'output' in outputs:
            # Limpa a resposta final
            resposta = outputs['output']
            resposta = resposta.replace("<think>", "").replace("</think>", "")
            resposta = re.sub(r">.*?<", "", resposta)
            resposta = re.sub(r"enviar_email\(.*?\)", "", resposta)  # Remove chamadas de email
            resposta = re.sub(r"\{input\}", "", resposta)  # Remove placeholders
            self.resposta_final = resposta.strip()
            self.em_pensamento = False
    
    def on_tool_start(self, *args, **kwargs):
        """Chamado quando uma ferramenta começa a ser usada."""
        tool_name = kwargs.get('tool', '')
        if tool_name != self.ultima_ferramenta:
            self.pensamentos.append(f"\n*Vou usar a função {tool_name} para isso...*")
            self.ultima_ferramenta = tool_name
    
    def on_tool_end(self, *args, **kwargs):
        """Chamado quando uma ferramenta termina de ser usada."""
        if 'output' in kwargs:
            self.pensamentos.append(f"\n*Resultado obtido: {kwargs['output']}*")
    
    def on_agent_action(self, *args, **kwargs):
        """Chamado quando o agente toma uma ação."""
        if 'tool' in kwargs and kwargs['tool'] != self.ultima_ferramenta:
            self.pensamentos.append(f"\n*Decidindo usar a função {kwargs['tool']}...*")
    
    def on_agent_finish(self, *args, **kwargs):
        """Chamado quando o agente termina."""
        self.em_pensamento = False

class CompreensaoMensagem:
    def __init__(self):
        # Intenções com contexto mais rico
        self.intencoes = {
            'consulta': {
                'palavras': ['qual', 'como', 'quando', 'onde', 'quem', 'por que', 'porque', 'temperatura', 'clima', 'hora', 'data'],
                'contextos': ['informação', 'dúvida', 'pergunta', 'consulta'],
                'ferramenta': None,
                'exemplos': ['qual é a temperatura hoje?', 'que horas são?', 'qual é a data?']
            },
            'listar': {
                'palavras': ['listar', 'mostrar', 'ver', 'arquivos', 'pastas', 'conteúdo', 'lista', 'diretório', 'pasta'],
                'contextos': ['arquivos', 'pastas', 'diretórios', 'conteúdo', 'estrutura'],
                'ferramenta': 'listar_arquivos',
                'exemplos': ['mostre os arquivos do desktop', 'liste o conteúdo da pasta documentos']
            },
            'criar': {
                'palavras': ['criar', 'novo', 'nova', 'fazer', 'adicionar', 'gerar', 'iniciar'],
                'contextos': ['arquivo', 'pasta', 'diretório', 'estrutura', 'projeto'],
                'ferramenta': 'criar_estrutura_pastas',
                'exemplos': ['crie uma pasta chamada projeto', 'faça um novo arquivo python']
            },
            'editar': {
                'palavras': ['editar', 'modificar', 'alterar', 'mudar', 'atualizar', 'ajustar'],
                'contextos': ['arquivo', 'código', 'texto', 'conteúdo'],
                'ferramenta': 'criar_arquivo_codigo',
                'exemplos': ['edite o arquivo main.py', 'modifique o código da função']
            },
            'executar': {
                'palavras': ['executar', 'rodar', 'iniciar', 'começar', 'rodar', 'testar'],
                'contextos': ['comando', 'programa', 'script', 'código'],
                'ferramenta': 'executar_comando',
                'exemplos': ['execute o comando dir', 'rode o script python']
            },
            'ajuda': {
                'palavras': ['ajuda', 'como', 'pode', 'quero', 'preciso', 'dúvida', 'explicar'],
                'contextos': ['ajuda', 'instrução', 'explicação', 'tutorial'],
                'ferramenta': None,
                'exemplos': ['como posso criar um arquivo?', 'me ajude a listar os arquivos']
            },
            'saudacao': {
                'palavras': ['oi', 'olá', 'bom dia', 'boa tarde', 'boa noite', 'hey', 'e aí'],
                'contextos': ['saudação', 'cumprimento', 'início'],
                'ferramenta': None,
                'exemplos': ['olá, tudo bem?', 'bom dia, como vai?']
            }
        }
        
        # Diretórios comuns com mais alternativas
        self.diretorios_comuns = {
            'desktop': ['desktop', 'área de trabalho', 'área', 'desktop', 'área de trabalho'],
            'documentos': ['documentos', 'documents', 'docs', 'documento'],
            'downloads': ['downloads', 'download', 'baixados', 'downloads'],
            'imagens': ['imagens', 'pictures', 'fotos', 'photos'],
            'musicas': ['músicas', 'music', 'sons', 'songs'],
            'videos': ['vídeos', 'videos', 'filmes', 'movies']
        }
        
        # Padrões de extração de parâmetros
        self.padroes = {
            'diretorio': r'(?:em|do|da|no|na|para|em|dentro de|os arquivos do|o conteúdo do)\s+([^\s]+(?:\s+[^\s]+)*)',
            'nome_arquivo': r'(?:arquivo|file)\s+(?:chamado|nomeado|com nome)\s+([^\s]+(?:\.[^\s]+)?)',
            'nome_pasta': r'(?:pasta|diretorio|folder)\s+(?:chamado|nomeado|com nome)\s+([^\s]+)',
            'comando': r'(?:comando|command)\s+(?:para|para executar|para rodar)\s+([^\s]+(?:\s+[^\s]+)*)'
        }
    
    def validar_mensagem(self, mensagem: str) -> tuple[bool, str]:
        """Validação mais robusta da mensagem."""
        mensagem = mensagem.strip()
        
        # Mensagem vazia
        if not mensagem:
            return False, "Por favor, digite algo para que eu possa ajudar."
        
        # Apenas pontuação ou caracteres especiais
        if re.match(r'^[\s\W]+$', mensagem):
            return False, "Não entendi o que você quer dizer. Poderia explicar melhor?"
        
        # Mensagem muito curta (menos de 2 caracteres)
        if len(mensagem) < 2:
            return False, "Sua mensagem está muito curta. Poderia fornecer mais detalhes?"
        
        # Muitos caracteres especiais
        if len(re.findall(r'[^\w\s]', mensagem)) > len(mensagem) * 0.5:
            return False, "Sua mensagem contém muitos caracteres especiais. Poderia reformular?"
        
        return True, ""
    
    def analisar_mensagem(self, mensagem: str) -> dict:
        """Análise mais robusta da mensagem."""
        # Validação básica
        valida, erro = self.validar_mensagem(mensagem)
        if not valida:
            return {
                'intencao': None,
                'ferramenta': None,
                'parametros': {},
                'confianca': 0.0,
                'erro': erro,
                'sugestoes': []
            }
        
        mensagem = mensagem.lower().strip()
        resultado = {
            'intencao': None,
            'ferramenta': None,
            'parametros': {},
            'confianca': 0.0,
            'erro': None,
            'sugestoes': []
        }
        
        # Verifica cada intenção
        for intencao, dados in self.intencoes.items():
            # Conta matches de palavras-chave
            matches_palavras = sum(1 for palavra in dados['palavras'] if palavra in mensagem)
            matches_contextos = sum(1 for contexto in dados['contextos'] if contexto in mensagem)
            
            if matches_palavras > 0 or matches_contextos > 0:
                # Calcula confiança baseada em múltiplos fatores
                confianca = 0.0
                
                # Base por ter encontrado palavras-chave
                if matches_palavras > 0:
                    confianca += 0.4
                    confianca += min(matches_palavras * 0.1, 0.3)  # Bônus por mais matches
                
                # Bônus por ter encontrado contexto
                if matches_contextos > 0:
                    confianca += 0.2
                    confianca += min(matches_contextos * 0.1, 0.2)  # Bônus por mais contextos
                
                # Bônus por similaridade com exemplos
                for exemplo in dados['exemplos']:
                    if self._calcular_similaridade(mensagem, exemplo) > 0.7:
                        confianca += 0.1
                
                # Limita a confiança
                confianca = min(confianca, 1.0)
                
                if confianca > resultado['confianca']:
                    resultado.update({
                        'intencao': intencao,
                        'ferramenta': dados['ferramenta'],
                        'confianca': confianca,
                        'sugestoes': dados['exemplos']
                    })
        
        # Se encontrou uma intenção, tenta extrair parâmetros
        if resultado['intencao']:
            self._extrair_parametros(mensagem, resultado)
        
        return resultado
    
    def _calcular_similaridade(self, texto1: str, texto2: str) -> float:
        """Calcula similaridade entre dois textos."""
        # Implementação simples usando palavras em comum
        palavras1 = set(texto1.lower().split())
        palavras2 = set(texto2.lower().split())
        
        if not palavras1 or not palavras2:
            return 0.0
        
        intersecao = len(palavras1.intersection(palavras2))
        uniao = len(palavras1.union(palavras2))
        
        return intersecao / uniao if uniao > 0 else 0.0
    
    def _extrair_parametros(self, mensagem: str, resultado: dict) -> None:
        """Extrai parâmetros da mensagem baseado na intenção."""
        # Para listar arquivos
        if resultado['intencao'] == 'listar':
            # Procura por diretório específico
            for dir_padrao, alternativas in self.diretorios_comuns.items():
                if any(alt in mensagem.lower() for alt in alternativas):
                    resultado['parametros']['diretorio'] = os.path.expanduser(f"~/{dir_padrao}")
                    resultado['confianca'] = min(resultado['confianca'] + 0.2, 1.0)
                    break
            
            # Tenta extrair diretório usando padrão
            match = re.search(self.padroes['diretorio'], mensagem.lower())
            if match:
                diretorio = match.group(1).strip()
                # Converte diretórios comuns para caminhos reais
                for dir_padrao, alternativas in self.diretorios_comuns.items():
                    if diretorio.lower() in alternativas:
                        diretorio = os.path.expanduser(f"~/{dir_padrao}")
                        break
                resultado['parametros']['diretorio'] = diretorio
                resultado['confianca'] = min(resultado['confianca'] + 0.1, 1.0)
        
        # Para criar arquivos/pastas
        elif resultado['intencao'] == 'criar':
            # Procura por nome de arquivo
            match = re.search(self.padroes['nome_arquivo'], mensagem)
            if match:
                nome = match.group(1).strip()
                resultado['parametros']['nome'] = nome
                resultado['confianca'] = min(resultado['confianca'] + 0.2, 1.0)
            
            # Procura por nome de pasta
            match = re.search(self.padroes['nome_pasta'], mensagem)
            if match:
                nome = match.group(1).strip()
                resultado['parametros']['nome'] = nome
                resultado['confianca'] = min(resultado['confianca'] + 0.2, 1.0)
        
        # Para executar comandos
        elif resultado['intencao'] == 'executar':
            match = re.search(self.padroes['comando'], mensagem)
            if match:
                comando = match.group(1).strip()
                resultado['parametros']['comando'] = comando
                resultado['confianca'] = min(resultado['confianca'] + 0.2, 1.0)

@dataclass
class ConfiguracaoAgente:
    """Configuração validada do agente."""
    modelo_geral: str
    modelo_coder: str
    temperatura: float
    top_p: float
    max_iteracoes: int
    timeout: int
    streaming: bool
    
    @classmethod
    def validar(cls, config: Dict[str, Any]) -> Resultado['ConfiguracaoAgente']:
        """Valida a configuração do agente."""
        try:
            # Validação dos modelos
            resultado_geral = validar_modelo(config.get('modelo_geral', 'qwen3-4b'))
            if resultado_geral.falha:
                return Resultado(erro=resultado_geral.erro)
                
            resultado_coder = validar_modelo(config.get('modelo_coder', 'llama-3.2-1b-instruct'))
            if resultado_coder.falha:
                return Resultado(erro=resultado_coder.erro)
            
            # Validação dos parâmetros numéricos
            temperatura = config.get('temperatura', 0.7)
            if not isinstance(temperatura, (int, float)) or not 0 <= temperatura <= 1:
                return Resultado(erro="Temperatura deve ser um número entre 0 e 1")
            
            top_p = config.get('top_p', 0.9)
            if not isinstance(top_p, (int, float)) or not 0 <= top_p <= 1:
                return Resultado(erro="Top_p deve ser um número entre 0 e 1")
            
            max_iteracoes = config.get('max_iteracoes', 5)
            if not isinstance(max_iteracoes, int) or max_iteracoes < 1:
                return Resultado(erro="Max_iteracoes deve ser um inteiro positivo")
            
            timeout = config.get('timeout', 30)
            if not isinstance(timeout, int) or timeout < 1:
                return Resultado(erro="Timeout deve ser um inteiro positivo")
            
            streaming = config.get('streaming', True)
            if not isinstance(streaming, bool):
                return Resultado(erro="Streaming deve ser um booleano")
            
            return Resultado(valor=cls(
                modelo_geral=resultado_geral.valor,
                modelo_coder=resultado_coder.valor,
                temperatura=temperatura,
                top_p=top_p,
                max_iteracoes=max_iteracoes,
                timeout=timeout,
                streaming=streaming
            ))
        except Exception as e:
            return Resultado(erro=f"Erro ao validar configuração: {str(e)}")

class FerramentaInvalidaError(Exception):
    """Exceção para ferramentas inválidas."""
    pass

def validar_ferramenta(ferramenta: BaseTool) -> Resultado[BaseTool]:
    """Valida uma ferramenta do agente."""
    try:
        if not isinstance(ferramenta, BaseTool):
            return Resultado(erro=f"Ferramenta deve ser uma instância de BaseTool: {type(ferramenta)}")
        
        if not ferramenta.name or not isinstance(ferramenta.name, str):
            return Resultado(erro="Ferramenta deve ter um nome válido")
        
        if not ferramenta.description or not isinstance(ferramenta.description, str):
            return Resultado(erro="Ferramenta deve ter uma descrição válida")
        
        if not callable(ferramenta._run):
            return Resultado(erro="Ferramenta deve ter um método _run executável")
        
        return Resultado(valor=ferramenta)
    except Exception as e:
        return Resultado(erro=f"Erro ao validar ferramenta: {str(e)}")

@dataclass
class EstadoAgente:
    """Estado validado do agente."""
    executores_inicializados: bool
    ferramentas_carregadas: bool
    historico_ativo: bool
    ultima_mensagem: Optional[Mensagem]
    ultimo_erro: Optional[str]
    nivel_seguranca: Optional['NivelSeguranca']
    
    @classmethod
    def inicial(cls) -> 'EstadoAgente':
        """Retorna o estado inicial do agente."""
        return cls(
            executores_inicializados=False,
            ferramentas_carregadas=False,
            historico_ativo=False,
            ultima_mensagem=None,
            ultimo_erro=None,
            nivel_seguranca=None
        )
    
    def validar_transicao(self, novo_estado: 'EstadoAgente') -> Resultado[None]:
        """Valida uma transição de estado."""
        try:
            # Validação de invariantes
            if novo_estado.executores_inicializados and not novo_estado.ferramentas_carregadas:
                return Resultado(erro="Executores não podem estar inicializados sem ferramentas carregadas")
            
            if novo_estado.ultima_mensagem and not isinstance(novo_estado.ultima_mensagem, Mensagem):
                return Resultado(erro="Última mensagem deve ser uma instância de Mensagem")
            
            if novo_estado.ultimo_erro and not isinstance(novo_estado.ultimo_erro, str):
                return Resultado(erro="Último erro deve ser uma string")
            
            return Resultado()
        except Exception as e:
            return Resultado(erro=f"Erro ao validar transição de estado: {str(e)}")

class NivelSeguranca(Enum):
    """Níveis de segurança para validação."""
    MINIMO = auto()
    MEDIO = auto()
    MAXIMO = auto()

@dataclass
class RequisitosAmbiente:
    """Requisitos do ambiente de execução."""
    python_version: Tuple[int, int, int]
    memoria_minima: int  # em MB
    cpu_minima: int  # em cores
    dependencias: Dict[str, str]
    permissoes_arquivo: Set[str]
    permissoes_rede: bool
    nivel_seguranca: NivelSeguranca

@dataclass
class StatusAmbiente:
    """Status validado do ambiente."""
    python_version: Tuple[int, int, int]
    memoria_disponivel: int
    cpu_disponivel: int
    dependencias_ok: bool
    permissoes_ok: bool
    rede_ok: bool
    ultima_verificacao: datetime
    erros: List[str] = field(default_factory=list)

class AmbienteError(Exception):
    """Exceção para erros de ambiente."""
    pass

class DependenciaError(Exception):
    """Exceção para erros de dependência."""
    pass

class SegurancaError(Exception):
    """Exceção para erros de segurança."""
    pass

def validar_ambiente(requisitos: RequisitosAmbiente) -> Resultado[StatusAmbiente]:
    """Valida o ambiente de execução."""
    try:
        status = StatusAmbiente(
            python_version=sys.version_info[:3],
            memoria_disponivel=psutil.virtual_memory().available // (1024 * 1024),
            cpu_disponivel=psutil.cpu_count(),
            dependencias_ok=True,
            permissoes_ok=True,
            rede_ok=True,
            ultima_verificacao=datetime.now(),
            erros=[]
        )
        
        # Validação da versão do Python
        if status.python_version < requisitos.python_version:
            status.erros.append(f"Versão do Python insuficiente: {status.python_version} < {requisitos.python_version}")
            status.dependencias_ok = False
        
        # Validação de memória
        if status.memoria_disponivel < requisitos.memoria_minima:
            status.erros.append(f"Memória insuficiente: {status.memoria_disponivel}MB < {requisitos.memoria_minima}MB")
            status.dependencias_ok = False
        
        # Validação de CPU
        if status.cpu_disponivel < requisitos.cpu_minima:
            status.erros.append(f"CPU insuficiente: {status.cpu_disponivel} cores < {requisitos.cpu_minima} cores")
            status.dependencias_ok = False
        
        # Validação de dependências
        for pacote, versao in requisitos.dependencias.items():
            try:
                pkg_resources.require(f"{pacote}=={versao}")
            except (pkg_resources.VersionConflict, pkg_resources.DistributionNotFound) as e:
                status.erros.append(f"Erro na dependência {pacote}: {str(e)}")
                status.dependencias_ok = False
        
        # Validação de permissões
        for permissao in requisitos.permissoes_arquivo:
            if not os.access(".", getattr(os, permissao.upper())):
                status.erros.append(f"Sem permissão {permissao}")
                status.permissoes_ok = False
        
        # Validação de rede
        if requisitos.permissoes_rede:
            try:
                socket.create_connection(("8.8.8.8", 53), timeout=3)
            except OSError:
                status.erros.append("Sem acesso à rede")
                status.rede_ok = False
        
        if status.erros:
            return Resultado(erro="\n".join(status.erros))
        
        return Resultado(valor=status)
    except Exception as e:
        return Resultado(erro=f"Erro ao validar ambiente: {str(e)}")

@contextmanager
def recuperacao_falha(max_tentativas: int = 3, delay: float = 1.0):
    """Contexto para recuperação de falhas."""
    tentativas = 0
    ultimo_erro = None
    
    while tentativas < max_tentativas:
        try:
            yield
            return
        except Exception as e:
            tentativas += 1
            ultimo_erro = e
            if tentativas < max_tentativas:
                logger.warning(f"Tentativa {tentativas} falhou: {str(e)}. Tentando novamente em {delay}s...")
                time.sleep(delay)
    
    raise AgenteError(f"Falha após {max_tentativas} tentativas. Último erro: {str(ultimo_erro)}")

def validar_seguranca(func):
    """Decorator para validação de segurança."""
    @wraps(func)
    def wrapper(self, *args, **kwargs):
        if self._estado.nivel_seguranca == NivelSeguranca.MAXIMO:
            # Validações extras de segurança
            if not self._validar_entrada_segura(*args, **kwargs):
                raise SegurancaError("Entrada não passou na validação de segurança")
            
            # Verifica assinatura digital se aplicável
            if hasattr(self, '_verificar_assinatura'):
                if not self._verificar_assinatura(*args, **kwargs):
                    raise SegurancaError("Assinatura inválida")
        
        return func(self, *args, **kwargs)
    return wrapper

class SegurancaAgente:
    """Classe para validações de segurança do agente."""
    
    @staticmethod
    def validar_entrada(entrada: str) -> Resultado[None]:
        """Valida entradas com foco em segurança."""
        try:
            # Validação de injeção de código
            if any(cmd in entrada.lower() for cmd in [
                "exec(", "eval(", "import ", "__import__", "os.system", "subprocess",
                "shell=True", "eval(", "compile(", "exec(", "__builtins__"
            ]):
                return Resultado(erro="Tentativa de injeção de código detectada")
            
            # Validação de tamanho
            if len(entrada) > 10000:  # Limite de 10KB
                return Resultado(erro="Entrada muito grande")
            
            # Validação de caracteres perigosos
            if any(c in entrada for c in ["\0", "\x1a", "\x00", "\x1b", "\x7f"]):
                return Resultado(erro="Caracteres perigosos detectados")
            
            # Validação de padrões suspeitos
            padroes_suspeitos = [
                r"\.\./",  # Path traversal
                r"<script>",  # XSS
                r"javascript:",  # XSS
                r"data:",  # Data URI
                r"vbscript:",  # VBScript
                r"on\w+\s*=",  # Event handlers
                r"document\.cookie",  # Cookie access
                r"localStorage",  # Local storage
                r"sessionStorage",  # Session storage
                r"indexedDB",  # IndexedDB
                r"fetch\(",  # Fetch API
                r"XMLHttpRequest",  # XHR
                r"WebSocket",  # WebSocket
                r"Worker",  # Web Worker
                r"SharedWorker"  # Shared Worker
            ]
            
            for padrao in padroes_suspeitos:
                if re.search(padrao, entrada, re.IGNORECASE):
                    return Resultado(erro=f"Padrão suspeito detectado: {padrao}")
            
            return Resultado()
            
        except Exception as e:
            return Resultado(erro=f"Erro na validação de segurança: {str(e)}")

    @staticmethod
    def validar_assinatura(dados: Any, assinatura: str) -> Resultado[None]:
        """Valida assinatura digital."""
        try:
            # Implementação da validação de assinatura
            # Por enquanto, apenas um placeholder
            return Resultado()
        except Exception as e:
            return Resultado(erro=f"Erro na validação de assinatura: {str(e)}")

    @staticmethod
    def validar_recursos() -> Resultado[None]:
        """Valida recursos do sistema."""
        try:
            processo = psutil.Process()
            
            # Validação de CPU
            if processo.cpu_percent() > 90:
                return Resultado(erro="Uso de CPU muito alto")
            
            # Validação de memória
            if processo.memory_percent() > 80:
                return Resultado(erro="Uso de memória muito alto")
            
            # Validação de arquivos abertos
            if len(processo.open_files()) > 100:
                return Resultado(erro="Muitos arquivos abertos")
            
            return Resultado()
            
        except Exception as e:
            return Resultado(erro=f"Erro na validação de recursos: {str(e)}")

@dataclass
class AjusteConfiguracao:
    """Configuração para ajustes automáticos."""
    intervalo_minutos: int = 60  # Intervalo entre ajustes
    intencoes_permitidas: Set[str] = field(default_factory=lambda: {'consulta', 'listar', 'criar'})
    max_palavras_por_ajuste: int = 5
    min_feedbacks_negativos: int = 3
    caminho_relatorio: str = "relatorios_ajuste"

@dataclass
class RelatorioAjuste:
    """Relatório de ajustes realizados."""
    timestamp: datetime
    palavras_adicionadas: Dict[str, List[str]]  # intenção -> palavras
    feedbacks_analisados: int
    intencoes_ajustadas: List[str]
    metricas: Dict[str, float]  # métricas de performance

class AjusteAutomatico:
    """Gerencia ajustes automáticos do agente."""
    
    def __init__(self, agente: 'AgenteIA', config: Optional[AjusteConfiguracao] = None):
        self.agente = agente
        self.config = config or AjusteConfiguracao()
        self.ultimo_ajuste: Optional[datetime] = None
        self.relatorios: List[RelatorioAjuste] = []
        
        # Garante que o diretório de relatórios existe
        os.makedirs(self.config.caminho_relatorio, exist_ok=True)
    
    def iniciar_ajustes_automaticos(self):
        """Inicia o processo de ajustes automáticos."""
        import threading
        import time
        
        def loop_ajustes():
            while True:
                self.realizar_ajuste()
                time.sleep(self.config.intervalo_minutos * 60)
        
        thread = threading.Thread(target=loop_ajustes, daemon=True)
        thread.start()
        logger.info("Ajustes automáticos iniciados")
    
    def realizar_ajuste(self) -> Optional[RelatorioAjuste]:
        """Realiza um ajuste baseado nos feedbacks recentes."""
        try:
            # Verifica se já passou tempo suficiente desde o último ajuste
            if self.ultimo_ajuste and (datetime.now() - self.ultimo_ajuste).total_seconds() < self.config.intervalo_minutos * 60:
                return None
            
            # Coleta feedbacks negativos recentes
            feedbacks_negativos = self._coletar_feedbacks_negativos()
            if len(feedbacks_negativos) < self.config.min_feedbacks_negativos:
                logger.info("Feedbacks insuficientes para ajuste")
                return None
            
            # Analisa os feedbacks
            palavras_por_intencao = self._analisar_feedbacks(feedbacks_negativos)
            
            # Cria relatório
            relatorio = RelatorioAjuste(
                timestamp=datetime.now(),
                palavras_adicionadas=palavras_por_intencao,
                feedbacks_analisados=len(feedbacks_negativos),
                intencoes_ajustadas=list(palavras_por_intencao.keys()),
                metricas=self._calcular_metricas()
            )
            
            # Aplica os ajustes
            self._aplicar_ajustes(palavras_por_intencao)
            
            # Salva o relatório
            self._salvar_relatorio(relatorio)
            
            self.ultimo_ajuste = datetime.now()
            self.relatorios.append(relatorio)
            
            logger.info(f"Ajuste realizado com sucesso: {len(palavras_por_intencao)} intenções ajustadas")
            return relatorio
            
        except Exception as e:
            logger.error(f"Erro ao realizar ajuste: {str(e)}")
            return None
    
    def _coletar_feedbacks_negativos(self) -> List[Dict[str, Any]]:
        """Coleta feedbacks negativos recentes."""
        feedbacks = []
        try:
            with open("feedbacks.jsonl", "r", encoding="utf-8") as f:
                for linha in f:
                    registro = json.loads(linha)
                    if not registro["util"]:
                        feedbacks.append(registro)
        except Exception as e:
            logger.error(f"Erro ao coletar feedbacks: {str(e)}")
        return feedbacks
    
    def _analisar_feedbacks(self, feedbacks: List[Dict[str, Any]]) -> Dict[str, List[str]]:
        """Analisa feedbacks e retorna palavras por intenção."""
        from collections import Counter
        
        palavras_por_intencao = {intencao: [] for intencao in self.config.intencoes_permitidas}
        
        for feedback in feedbacks:
            prompt = feedback["prompt"].lower()
            palavras = prompt.split()
            
            # Conta frequência das palavras
            contador = Counter(palavras)
            
            # Para cada intenção permitida
            for intencao in self.config.intencoes_permitidas:
                # Verifica se o prompt se relaciona com a intenção
                if any(palavra in self.agente.compreensao.intencoes[intencao]['palavras'] for palavra in palavras):
                    # Adiciona as palavras mais frequentes que não estão nas palavras-chave
                    palavras_existentes = set(self.agente.compreensao.intencoes[intencao]['palavras'])
                    novas_palavras = [p for p, _ in contador.most_common(self.config.max_palavras_por_ajuste)
                                    if p not in palavras_existentes]
                    palavras_por_intencao[intencao].extend(novas_palavras)
        
        return palavras_por_intencao
    
    def _aplicar_ajustes(self, palavras_por_intencao: Dict[str, List[str]]):
        """Aplica os ajustes nas intenções."""
        for intencao, palavras in palavras_por_intencao.items():
            if intencao in self.agente.compreensao.intencoes:
                for palavra in palavras:
                    if palavra not in self.agente.compreensao.intencoes[intencao]['palavras']:
                        self.agente.compreensao.intencoes[intencao]['palavras'].append(palavra)
    
    def _calcular_metricas(self) -> Dict[str, float]:
        """Calcula métricas de performance do agente."""
        try:
            total_feedbacks = 0
            feedbacks_positivos = 0
            
            with open("feedbacks.jsonl", "r", encoding="utf-8") as f:
                for linha in f:
                    registro = json.loads(linha)
                    total_feedbacks += 1
                    if registro["util"]:
                        feedbacks_positivos += 1
            
            taxa_acerto = feedbacks_positivos / total_feedbacks if total_feedbacks > 0 else 0
            
            return {
                "taxa_acerto": taxa_acerto,
                "total_feedbacks": total_feedbacks,
                "feedbacks_positivos": feedbacks_positivos
            }
        except Exception as e:
            logger.error(f"Erro ao calcular métricas: {str(e)}")
            return {}
    
    def _salvar_relatorio(self, relatorio: RelatorioAjuste):
        """Salva o relatório em arquivo."""
        try:
            timestamp = relatorio.timestamp.strftime("%Y%m%d_%H%M%S")
            caminho = os.path.join(self.config.caminho_relatorio, f"ajuste_{timestamp}.json")
            
            with open(caminho, "w", encoding="utf-8") as f:
                json.dump({
                    "timestamp": relatorio.timestamp.isoformat(),
                    "palavras_adicionadas": relatorio.palavras_adicionadas,
                    "feedbacks_analisados": relatorio.feedbacks_analisados,
                    "intencoes_ajustadas": relatorio.intencoes_ajustadas,
                    "metricas": relatorio.metricas
                }, f, ensure_ascii=False, indent=2)
            
            logger.info(f"Relatório salvo em: {caminho}")
        except Exception as e:
            logger.error(f"Erro ao salvar relatório: {str(e)}")
    
    def listar_relatorios(self) -> List[str]:
        """Lista todos os relatórios disponíveis."""
        try:
            return [f for f in os.listdir(self.config.caminho_relatorio) 
                   if f.startswith("ajuste_") and f.endswith(".json")]
        except Exception as e:
            logger.error(f"Erro ao listar relatórios: {str(e)}")
            return []

class AgenteIA:
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Inicializa o agente com validação completa."""
        # Validação do ambiente
        requisitos = RequisitosAmbiente(
            python_version=(3, 8, 0),
            memoria_minima=1024,  # 1GB
            cpu_minima=2,
            dependencias={
                "langchain": "0.1.0",
                "openai": "1.0.0",
                "psutil": "5.9.0"
            },
            permissoes_arquivo={"R_OK", "W_OK"},
            permissoes_rede=True,
            nivel_seguranca=NivelSeguranca.MAXIMO
        )
        
        resultado_ambiente = validar_ambiente(requisitos)
        if resultado_ambiente.falha:
            raise AmbienteError(resultado_ambiente.erro)
        
        self._status_ambiente = resultado_ambiente.valor
        self._seguranca = SegurancaAgente()
        
        # Inicializa o estado com nível de segurança
        self._estado = EstadoAgente.inicial()
        self._estado.nivel_seguranca = requisitos.nivel_seguranca
        
        # Configuração de recuperação
        self._max_tentativas = 3
        self._delay_recuperacao = 1.0
        
        # Inicialização com recuperação
        with recuperacao_falha(self._max_tentativas, self._delay_recuperacao):
            # Validação da configuração
            resultado_config = ConfiguracaoAgente.validar(config or {})
            if resultado_config.falha:
                self._atualizar_estado(ultimo_erro=resultado_config.erro)
                raise ValidacaoError(resultado_config.erro)
            
            self.config = resultado_config.valor
            
            # Validação do diretório de histórico
            resultado_historico = validar_caminho("historico")
            if resultado_historico.falha:
                os.makedirs("historico", exist_ok=True)
            
            # Inicialização dos componentes com validação
            self.compreensao = CompreensaoMensagem()
            self.contexto_manager = ContextoManager()
            self.callback_handler = PensamentoCallback()
            
            # Validação da mensagem do sistema
            if not isinstance(self.system_message, str):
                self._atualizar_estado(ultimo_erro="Mensagem do sistema deve ser uma string")
                raise ValidacaoError("Mensagem do sistema deve ser uma string")
            
            # Inicializa os executores com validação
            self._inicializar_executores()
            
            # Atualiza o estado após inicialização bem-sucedida
            self._atualizar_estado(
                executores_inicializados=True,
                ferramentas_carregadas=True,
                historico_ativo=True
            )
            
            logger.info(f"Agente inicializado com configuração: {self.config}")

            # Inicializa o sistema de ajustes automáticos
            self.ajuste_automatico = AjusteAutomatico(self)
            
            # Inicia os ajustes automáticos
            self.ajuste_automatico.iniciar_ajustes_automaticos()

            # Inicializa o ChatOllama para o modelo coder
            self.ollama = ChatOllama(
                model=self.config.modelo_geral,
                base_url=CONFIG["ollama"]["base_url"],
                request_timeout=CONFIG["ollama"]["timeout"]
            )

            self.ollama_coder = ChatOllama(
                model=self.config.modelo_coder,
                base_url=CONFIG["ollama"]["base_url"],
                request_timeout=CONFIG["ollama"]["timeout"]
            )

            self.provedor_atual = None
            self.llm_geral = None
            self.llm_coder = None
            self.ferramentas = []
            self._inicializar_provedor()

    def _inicializar_provedor(self) -> None:
        """Inicializa o provedor de modelo mais adequado."""
        if verificar_conectividade_openrouter():
            self._configurar_openrouter()
            self.provedor_atual = ProvedorModelo.OPENROUTER
        elif verificar_conectividade_ollama():
            self._configurar_ollama()
            self.provedor_atual = ProvedorModelo.OLLAMA
        else:
            raise AgenteError("Nenhum provedor de modelo disponível")

    def _configurar_openrouter(self) -> None:
        """Configura o provedor de modelo OpenRouter usando a classe customizada ChatOpenRouter."""
        self.logger.info("Iniciando configuração do OpenRouter usando a classe ChatOpenRouter customizada...")
        try:
            # Obter configuração específica do OpenRouter
            openrouter_config = self.config.get('provedores', {}).get('openrouter', {})
            api_key = openrouter_config.get('api_key') or os.getenv('OPENROUTER_API_KEY')

            if not api_key:
                 raise ValueError("A chave da API do OpenRouter não foi fornecida.")

            base_url = openrouter_config.get('base_url', 'https://openrouter.ai/api/v1')

            # Usando a classe ChatOpenRouter customizada
            self.llm_geral = ChatOpenRouter(
                model=openrouter_config.get('modelo_geral', 'meta-llama/llama-3.3-70b-instruct:free'),
                openai_api_key=api_key,
                base_url=base_url,
                temperature=openrouter_config.get('temperatura', 0.3),
                max_tokens=openrouter_config.get('max_tokens', 2000),
                timeout=openrouter_config.get('timeout', 30),
                # Adicionar headers customizados se necessário para identificação
                headers={
                    "HTTP-Referer": "https://github.com/yourusername/agenteia", # Substitua com o URL do seu projeto
                    "X-Title": "AgenteIA", # Substitua com o nome do seu app
                    "X-Provider": "openrouter"
                }
            )
            self.logger.info("Modelo geral configurado com sucesso usando ChatOpenRouter customizada")

            self.llm_coder = ChatOpenRouter(
                model=openrouter_config.get('modelo_coder', 'qwen/qwen-2.5-coder-32b-instruct:free'),
                openai_api_key=api_key,
                base_url=base_url,
                temperature=openrouter_config.get('temperatura', 0.3),
                max_tokens=openrouter_config.get('max_tokens', 2000),
                timeout=openrouter_config.get('timeout', 30),
                 headers={
                    "HTTP-Referer": "https://github.com/yourusername/agenteia", # Substitua com o URL do seu projeto
                    "X-Title": "AgenteIA", # Substitua com o nome do seu app
                    "X-Provider": "openrouter"
                }
            )
            self.logger.info("Modelo coder configurado com sucesso usando ChatOpenRouter customizada")

            self.provedor_ativo = ProvedorModelo.OPENROUTER
            self.logger.info("Modelos OpenRouter configurados com sucesso usando ChatOpenRouter customizada")
            self._atualizar_estado(executores_inicializados=False) # Forçar reinicialização dos executores

        except Exception as e:
            self.logger.error(f"Erro ao configurar provedor OpenRouter: {e}")
            # Opcional: Voltar para Ollama ou levantar a exceção
            # self._configurar_ollama() # Exemplo de fallback
            raise AgenteError(f"Falha ao configurar provedor OpenRouter: {e}") from e

    def _configurar_ollama(self) -> None:
        """Configura o Ollama como provedor."""
        try:
            from langchain_community.chat_models import ChatOllama
            config_ollama = self.config.get('llm', {})
            self.llm_geral = ChatOllama(
                base_url=config_ollama.get('url', 'http://localhost:11434'),
                model=config_ollama.get('model', 'qwen3-1.7b'),
                temperature=config_ollama.get('temperature', 0.3),
                timeout=config_ollama.get('timeout', 30)
            )
            self.llm_coder = ChatOllama(
                base_url=config_ollama.get('url', 'http://localhost:11434'),
                model=config_ollama.get('coder', 'qwen2.5-coder:3b'),
                temperature=config_ollama.get('temperature', 0.3),
                timeout=config_ollama.get('timeout', 30)
            )
        except Exception as e:
            raise AgenteError(f"Erro ao configurar Ollama: {str(e)}")

    def alternar_provedor(self) -> None:
        """Alterna entre provedores de modelo."""
        if self.provedor_atual == ProvedorModelo.OPENROUTER:
            if verificar_conectividade_ollama():
                self._configurar_ollama()
                self.provedor_atual = ProvedorModelo.OLLAMA
            else:
                raise AgenteError("Ollama não está disponível")
        else:
            if verificar_conectividade_openrouter():
                self._configurar_openrouter()
                self.provedor_atual = ProvedorModelo.OPENROUTER
            else:
                raise AgenteError("OpenRouter não está disponível")

    def _atualizar_estado(self, **kwargs) -> None:
        """Atualiza o estado do agente com validação."""
        try:
            novo_estado = EstadoAgente(
                executores_inicializados=kwargs.get('executores_inicializados', self._estado.executores_inicializados),
                ferramentas_carregadas=kwargs.get('ferramentas_carregadas', self._estado.ferramentas_carregadas),
                historico_ativo=kwargs.get('historico_ativo', self._estado.historico_ativo),
                ultima_mensagem=kwargs.get('ultima_mensagem', self._estado.ultima_mensagem),
                ultimo_erro=kwargs.get('ultimo_erro', self._estado.ultimo_erro),
                nivel_seguranca=kwargs.get('nivel_seguranca', self._estado.nivel_seguranca)
            )
            
            resultado = self._estado.validar_transicao(novo_estado)
            if resultado.falha:
                raise ValidacaoError(resultado.erro)
            
            self._estado = novo_estado
            logger.debug(f"Estado do agente atualizado: {self._estado}")
            
        except Exception as e:
            logger.error(f"Erro ao atualizar estado: {str(e)}")
            raise AgenteError(f"Erro ao atualizar estado: {str(e)}")

    def _validar_ferramentas(self, ferramentas: List[BaseTool]) -> Resultado[List[BaseTool]]:
        """Valida uma lista de ferramentas."""
        ferramentas_validadas = []
        for ferramenta in ferramentas:
            resultado = validar_ferramenta(ferramenta)
            if resultado.falha:
                return Resultado(erro=resultado.erro)
            ferramentas_validadas.append(resultado.valor)
        return Resultado(valor=ferramentas_validadas)

    def _inicializar_executores(self) -> None:
        """Inicializa os executores do agente com validação."""
        try:
            # Obter e validar ferramentas
            ferramentas_gerais = self._obter_ferramentas_gerais()
            resultado_gerais = self._validar_ferramentas(ferramentas_gerais)
            if resultado_gerais.falha:
                raise FerramentaInvalidaError(resultado_gerais.erro)
            
            ferramentas_codigo = self._obter_ferramentas_codigo()
            resultado_codigo = self._validar_ferramentas(ferramentas_codigo)
            if resultado_codigo.falha:
                raise FerramentaInvalidaError(resultado_codigo.erro)
            
            # Criar executores com ferramentas validadas
            self.agent_executor_geral = self._criar_agente(
                ferramentas=resultado_gerais.valor,
                modelo=self.config.modelo_geral
            )
            logger.info("Executor geral inicializado com sucesso")
            
            self.agent_executor_coder = self._criar_agente(
                ferramentas=resultado_codigo.valor,
                modelo=self.config.modelo_coder
            )
            logger.info("Executor de código inicializado com sucesso")
            
        except Exception as e:
            logger.error(f"Erro ao inicializar executores: {str(e)}")
            raise AgenteError(f"Erro ao inicializar executores: {str(e)}")
    
    def _criar_agente(self, ferramentas: List[BaseTool], modelo: str) -> AgentExecutor:
        """Cria um agente com validação rigorosa."""
        try:
            # Validar ferramentas
            resultado = self._validar_ferramentas(ferramentas)
            if resultado.falha:
                raise FerramentaInvalidaError(resultado.erro)
            
            # Configurar o modelo com parâmetros validados
            llm_stream = ChatOllama(
                model=modelo,
                base_url=CONFIG["ollama"]["base_url"],
                request_timeout=CONFIG["ollama"]["timeout"],
                streaming=True,
                temperature=CONFIG["modelo"]["temperature"],
                top_p=CONFIG["modelo"]["top_p"],
                max_tokens=CONFIG["modelo"]["max_tokens"]
            )

            # Criar o agente com validação
            agent = create_tool_calling_agent(
                llm=llm_stream,
                tools=resultado.valor,
                system_message=SystemMessage(content=self.system_message),
                verbose=True
            )

            # Criar o executor com validação
            executor = AgentExecutor(
                agent=agent,
                tools=resultado.valor,
                verbose=True,
                handle_parsing_errors=True,
                max_iterations=self.config.max_iteracoes,
                callbacks=[self.callback_handler]
            )

            return executor

        except Exception as e:
            logger.error(f"Erro ao criar agente: {e}")
            raise AgenteError(f"Erro ao criar agente: {e}")

    @validar_seguranca
    def processar_mensagem(self, mensagem: str, mostrar_raciocinio: bool = True, 
                          temperatura: Optional[float] = None, 
                          top_p: Optional[float] = None, 
                          perfil: Optional[str] = None) -> Resultado[str]:
        """Processa uma mensagem com validação de segurança."""
        try:
            # Validação do estado
            if not self._estado.executores_inicializados:
                return Resultado(erro="Agente não está inicializado corretamente")
            
            if not self._estado.ferramentas_carregadas:
                return Resultado(erro="Ferramentas não estão carregadas")
            
            # Validação da mensagem
            if not isinstance(mensagem, str):
                self._atualizar_estado(ultimo_erro="Mensagem deve ser uma string")
                return Resultado(erro="Mensagem deve ser uma string")
            
            mensagem = mensagem.strip()
            if not mensagem:
                self._atualizar_estado(ultimo_erro="Mensagem não pode ser vazia")
                return Resultado(erro="Mensagem não pode ser vazia")
            
            # Validação de segurança
            resultado_seguranca = self._seguranca.validar_entrada(mensagem)
            if resultado_seguranca.falha:
                self._atualizar_estado(ultimo_erro=resultado_seguranca.erro)
                return Resultado(erro=resultado_seguranca.erro)
            
            # Validação de recursos
            resultado_recursos = self._seguranca.validar_recursos()
            if resultado_recursos.falha:
                logger.warning(resultado_recursos.erro)
            
            # Processamento com recuperação
            with recuperacao_falha(self._max_tentativas, self._delay_recuperacao):
                # Cria e valida a mensagem
                try:
                    msg = Mensagem.criar(TipoMensagem.USUARIO, mensagem)
                    self._atualizar_estado(ultima_mensagem=msg)
                except ValueError as e:
                    self._atualizar_estado(ultimo_erro=str(e))
                    return Resultado(erro=str(e))
                
                # Processa a mensagem
                resultado = self._processar_mensagem_interna(msg, mostrar_raciocinio, temperatura, top_p, perfil)
                
                # Atualiza o estado
                if resultado.sucesso:
                    self._atualizar_estado(ultimo_erro=None)
                else: # Adicionado else para cobrir o caso de falha
                    self._atualizar_estado(ultimo_erro=resultado.erro)
                
                return resultado
                
        except Exception as e:
            erro = f"Erro ao processar mensagem: {str(e)}"
            self._atualizar_estado(ultimo_erro=erro)
            logger.error(erro)
            return Resultado(erro=erro)

    def _escolher_executor(self, analise: dict) -> AgentExecutor:
        """Escolhe o executor apropriado baseado na análise."""
        if analise['intencao'] in ['editar', 'executar'] or 'código' in analise.get('parametros', {}):
            logger.info("Escolhendo executor coder para tarefa de programação")
            return self.agent_executor_coder
        else:
            logger.info("Escolhendo executor geral para tarefa não específica")
            return self.agent_executor_geral
    
    def _limpar_resposta(self, texto: str) -> str:
        """Limpa a resposta removendo tags e formatação indesejada."""
        # Remove tags HTML e markdown
        texto = re.sub(r'<[^>]+>', '', texto)
        texto = re.sub(r'\[.*?\]', '', texto)
        texto = re.sub(r'\*\*.*?\*\*', '', texto)
        texto = re.sub(r'\*.*?\*', '', texto)
        
        # Remove emojis e caracteres especiais
        texto = re.sub(r'[^\w\s.,!?-]', '', texto)
        
        # Remove espaços extras
        texto = ' '.join(texto.split())
        
        return texto.strip()
    
    def _formatar_resposta(self, pensamentos: str, resposta: str, mostrar_raciocinio: bool) -> str:
        """Formata a resposta final."""
        if not mostrar_raciocinio or not pensamentos:
            return resposta
        
        return f"🤔 **Pensando...**\n{pensamentos}\n\n💭 **Resposta:**\n{resposta}"
    
    def _gerar_mensagem_ajuda(self, sugestoes: list) -> str:
        """Gera uma mensagem de ajuda com sugestões."""
        if not sugestoes:
            sugestoes = [
                "mostre os arquivos do desktop",
                "crie uma pasta chamada projeto",
                "edite o arquivo main.py",
                "execute o comando dir"
            ]
        
        mensagem = "Claro! Posso ajudar você a:\n"
        mensagem += "- Listar arquivos e pastas\n"
        mensagem += "- Criar novos arquivos e pastas\n"
        mensagem += "- Editar arquivos\n"
        mensagem += "- Executar comandos\n\n"
        mensagem += "Alguns exemplos:\n"
        for sugestao in sugestoes[:3]:  # Limita a 3 sugestões
            mensagem += f"- {sugestao}\n"
        mensagem += "\nO que você gostaria de fazer?"
        
        return mensagem
    
    def _gerar_mensagem_confusao(self, sugestoes: list) -> str:
        """Gera uma mensagem quando não entendeu a solicitação."""
        if not sugestoes:
            sugestoes = [
                "mostre os arquivos do desktop",
                "crie uma pasta chamada projeto",
                "edite o arquivo main.py",
                "execute o comando dir"
            ]
        
        mensagem = "Não entendi bem o que você precisa. Você pode:\n"
        for sugestao in sugestoes[:3]:  # Limita a 3 sugestões
            mensagem += f"- {sugestao}\n"
        mensagem += "\nPoderia reformular sua solicitação?"
        
        return mensagem

    def _obter_ferramentas_gerais(self) -> List[Any]:
        """Obtém as ferramentas gerais do agente."""
        try:
            from ferramentas import (
                listar_arquivos,
                criar_estrutura_pastas,
                pesquisar_web,
                calcular,
                executar_comando
            )
            
            return [
                listar_arquivos,
                criar_estrutura_pastas,
                pesquisar_web,
                calcular,
                executar_comando
            ]
        except Exception as e:
            logger.error(f"Erro ao obter ferramentas gerais: {str(e)}")
            raise AgenteError(f"Erro ao obter ferramentas gerais: {str(e)}")

    def _obter_ferramentas_codigo(self) -> List[Any]:
        """Obtém as ferramentas específicas para código."""
        try:
            from ferramentas import (
                criar_arquivo_codigo,
                ler_arquivo,
                escrever_arquivo,
                criar_word,
                criar_excel,
                criar_ppt
            )
            
            return [
                criar_arquivo_codigo,
                ler_arquivo,
                escrever_arquivo,
                criar_word,
                criar_excel,
                criar_ppt
            ]
        except Exception as e:
            logger.error(f"Erro ao obter ferramentas de código: {str(e)}")
            raise AgenteError(f"Erro ao obter ferramentas de código: {str(e)}")

    def limpar_historico(self) -> Resultado[None]:
        """Limpa o histórico com validação de estado."""
        try:
            if not self._estado.historico_ativo:
                return Resultado(erro="Histórico não está ativo")
            
            self.contexto_manager.limpar_historico()
            self._atualizar_estado(ultima_mensagem=None, ultimo_erro=None)
            logger.info("Histórico limpo com sucesso")
            return Resultado()
            
        except Exception as e:
            erro = f"Erro ao limpar histórico: {str(e)}"
            self._atualizar_estado(ultimo_erro=erro)
            logger.error(erro)
            return Resultado(erro=erro)
            
    def salvar_historico(self, caminho: Optional[str] = None) -> Resultado[None]:
        """Salva o histórico com validação rigorosa."""
        try:
            if caminho is not None:
                resultado = validar_caminho(caminho)
                if resultado.falha:
                    return Resultado(erro=resultado.erro)
                caminho = str(resultado.valor)
            else:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                caminho = os.path.join('historico', f'historico_{timestamp}.json')
            
            # Validação do histórico
            historico = self.contexto_manager.obter_historico()
            if not isinstance(historico, list):
                return Resultado(erro="Histórico inválido")
            
            # Garante que o diretório existe
            os.makedirs(os.path.dirname(caminho), exist_ok=True)
            
            with open(caminho, 'w', encoding='utf-8') as f:
                json.dump(historico, f, ensure_ascii=False, indent=2)
                
            logger.info(f"Histórico salvo com sucesso em '{caminho}'")
            return Resultado()
            
        except Exception as e:
            erro = f"Erro ao salvar histórico: {str(e)}"
            logger.error(erro)
            return Resultado(erro=erro)
            
    def carregar_historico(self, caminho: str = None):
        """
        Carrega o histórico de mensagens de um arquivo JSON.
        
        Args:
            caminho (str, optional): Caminho do arquivo para carregar o histórico.
                                   Se não fornecido, carrega o arquivo mais recente do diretório historico.
        """
        try:
            if caminho is None:
                # Lista todos os arquivos de histórico
                arquivos = [f for f in os.listdir('historico') if f.startswith('historico_') and f.endswith('.json')]
                if not arquivos:
                    logger.warning("Nenhum arquivo de histórico encontrado")
                    return
                
                # Pega o arquivo mais recente
                arquivo_mais_recente = max(arquivos, key=lambda x: os.path.getctime(os.path.join('historico', x)))
                caminho = os.path.join('historico', arquivo_mais_recente)
            
            if os.path.exists(caminho):
                with open(caminho, 'r', encoding='utf-8') as f:
                    self.contexto_manager.historico = json.load(f)
                    
                logger.info(f"Histórico carregado com sucesso de '{caminho}'")
                
        except Exception as e:
            erro = f"Erro ao carregar histórico: {str(e)}"
            logger.error(erro)
            return erro

    def listar_historicos(self) -> List[str]:
        """
        Lista todos os arquivos de histórico disponíveis.
        
        Returns:
            List[str]: Lista com os nomes dos arquivos de histórico
        """
        try:
            if not os.path.exists('historico'):
                return []
            
            arquivos = [f for f in os.listdir('historico') if f.startswith('historico_') and f.endswith('.json')]
            arquivos.sort(reverse=True)  # Ordena do mais recente para o mais antigo
            return arquivos
        except Exception as e:
            logger.error(f"Erro ao listar históricos: {str(e)}")
            return []

    def processar_arquivo(self, caminho_arquivo: str) -> str:
        """Processa um arquivo enviado pelo usuário."""
        try:
            extensao = os.path.splitext(caminho_arquivo)[1].lower()
            
            if extensao in ['.txt', '.md', '.py', '.js', '.html', '.css']:
                with open(caminho_arquivo, 'r', encoding='utf-8') as f:
                    conteudo = f.read()
                return self.processar_mensagem(f"Analise este arquivo {extensao}:\n\n{conteudo}")
            
            elif extensao in ['.docx', '.doc']:
                from docx import Document
                doc = Document(caminho_arquivo)
                texto = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                return self.processar_mensagem(f"Analise este documento Word:\n\n{texto}")
            
            elif extensao in ['.xlsx', '.xls']:
                import pandas as pd
                df = pd.read_excel(caminho_arquivo)
                return self.processar_mensagem(f"Analise esta planilha Excel:\n\n{df.to_string()}")
            
            elif extensao in ['.pptx', '.ppt']:
                from pptx import Presentation
                prs = Presentation(caminho_arquivo)
                texto = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texto.append(shape.text)
                return self.processar_mensagem(f"Analise esta apresentação PowerPoint:\n\n{'\n'.join(texto)}")
            
            else:
                return f"Formato de arquivo não suportado: {extensao}"
                
        except Exception as e:
            erro = f"Erro ao processar arquivo: {str(e)}"
            logger.error(erro)
            return erro

    def solicitar_feedback(self, resposta: str) -> bool:
        """Solicita feedback do usuário sobre a resposta."""
        print("A resposta foi útil? (s/n): ")
        feedback = input().strip().lower()
        return feedback == 's'

    def registrar_feedback(self, prompt: str, resposta: str, util: bool):
        """Registra o feedback do usuário em um arquivo JSONL."""
        registro = {
            "prompt": prompt,
            "resposta": resposta,
            "util": util,
            "timestamp": datetime.now().isoformat()
        }
        try:
            with open("feedbacks.jsonl", "a", encoding="utf-8") as f:
                f.write(json.dumps(registro, ensure_ascii=False) + "\n")
            logger.info(f"Feedback registrado: {registro}")
        except Exception as e:
            logger.error(f"Erro ao registrar feedback: {str(e)}")

    def ajustar_regras_com_base_no_feedback(self, caminho: str = "feedbacks.jsonl"):
        """Ajusta regras de compreensão automaticamente com base nos feedbacks negativos."""
        from collections import Counter
        prompts_negativos = []
        try:
            with open(caminho, "r", encoding="utf-8") as f:
                for linha in f:
                    registro = json.loads(linha)
                    if not registro["util"]:
                        prompts_negativos.append(registro["prompt"])
            # Extrai palavras mais comuns dos prompts negativos
            palavras = []
            for prompt in prompts_negativos:
                palavras.extend(prompt.lower().split())
            mais_comuns = [p for p, _ in Counter(palavras).most_common(5)]
            print(f"Palavras mais comuns em feedbacks negativos: {mais_comuns}")
            # Adiciona essas palavras como palavras-chave extras para intenções existentes
            for palavra in mais_comuns:
                for intencao, dados in self.compreensao.intencoes.items():
                    if palavra not in dados['palavras']:
                        dados['palavras'].append(palavra)
            print("Regras de compreensão ajustadas automaticamente com base nos feedbacks negativos.")
        except Exception as e:
            print(f"Erro ao ajustar regras: {e}")

def verificar_conectividade_openrouter() -> bool:
    """Verifica se é possível conectar ao OpenRouter."""
    try:
        import requests
        response = requests.get("https://openrouter.ai/api/v1/health", timeout=5)
        return response.status_code == 200
    except:
        return False

def verificar_conectividade_ollama() -> bool:
    """Verifica se o Ollama está rodando localmente."""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/health", timeout=5)
        return response.status_code == 200
    except:
        return False

class ProvedorModelo(Enum):
    """Provedores de modelo disponíveis."""
    OLLAMA = "ollama"
    OPENROUTER = "openrouter"

class ChatOpenRouter(ChatOpenAI):
    """
    ChatOpenAI customizada para usar a API do OpenRouter, compatível com bind_tools.
    """
    openai_api_key: Optional[SecretStr] = Field(
        alias="api_key", default_factory=secret_from_env("OPENROUTER_API_KEY", default=None)
    )

    @property
    def lc_secrets(self) -> dict[str, str]:
        return {"openai_api_key": "OPENROUTER_API_KEY"}

    def __init__(self,
                 model: str, # Alterado de model_name para model para consistência com ChatOpenAI
                 openai_api_key: Optional[str] = None,
                 base_url: str = "https://openrouter.ai/api/v1", # Alterado de openai_api_base para base_url
                 **kwargs):
        openai_api_key = openai_api_key or os.environ.get("OPENROUTER_API_KEY")
        super().__init__(base_url=base_url, openai_api_key=openai_api_key, model=model, **kwargs) # Passa model como 'model'
