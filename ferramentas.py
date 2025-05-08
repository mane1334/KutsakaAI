"""
Módulo de ferramentas do Agente IA
Contém todas as ferramentas disponíveis para o agente, com melhorias de segurança e performance
"""

import os
import subprocess
import pyautogui
import time
import requests
import shutil
from bs4 import BeautifulSoup
from urllib.parse import unquote, parse_qs
from langchain_core.tools import tool
from functools import lru_cache
import logging
import math
from typing import Optional, List, Dict, Any
import json
from datetime import datetime
import config
from docx import Document
import openpyxl
from pptx import Presentation

# Configurações de segurança
MAX_FILE_SIZE = config.MAX_FILE_SIZE
ALLOWED_COMMANDS = config.ALLOWED_COMMANDS
ALLOWED_FILE_EXTENSIONS = config.ALLOWED_FILE_EXTENSIONS
REQUEST_TIMEOUT = config.REQUEST_TIMEOUT

class SecurityError(Exception):
    """Exceção para erros de segurança."""
    pass

def validar_caminho(caminho: str) -> bool:
    # Para acesso total, descomente a linha abaixo:
    # return True
    # Para acesso restrito, use a lista de diretórios permitidos:
    caminho_absoluto = os.path.abspath(caminho)
    return any(caminho_absoluto.startswith(os.path.abspath(d)) for d in config.ALLOWED_DIRS)

def validar_tamanho_arquivo(caminho: str) -> bool:
    """Valida se o tamanho do arquivo está dentro do limite permitido."""
    return os.path.getsize(caminho) <= MAX_FILE_SIZE

def validar_extensao(caminho: str) -> bool:
    """Valida se a extensão do arquivo é permitida."""
    _, ext = os.path.splitext(caminho)
    return ext.lower() in ALLOWED_FILE_EXTENSIONS

def formatar_resposta(dados: Dict[str, Any]) -> str:
    """Formata a resposta em JSON com indentação."""
    return json.dumps(dados, ensure_ascii=False, indent=2)

@tool
def pesquisar_web(query: str) -> str:
    """
    Realiza uma pesquisa na web usando DuckDuckGo e retorna um resumo dos primeiros resultados.
    Útil para encontrar informações atuais, fatos, ou quando o conhecimento interno do agente não é suficiente.
    """
    try:
        logging.info(f"Pesquisando na web por: '{query}'")
        url = f"https://html.duckduckgo.com/html/?q={query}"
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, headers=headers, timeout=REQUEST_TIMEOUT)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.text, 'html.parser')
        resultados = []
        
        for result_div in soup.find_all('div', class_='web-result', limit=3):
            title_tag = result_div.find('a', class_='result__a')
            snippet_tag = result_div.find('a', class_='result__snippet')
            
            if title_tag and snippet_tag:
                title = title_tag.get_text(strip=True)
                snippet = snippet_tag.get_text(strip=True)
                raw_link = title_tag['href']
                link = raw_link
                
                if "duckduckgo.com/l/" in raw_link and "uddg=" in raw_link:
                    try:
                        params_str = raw_link.split('uddg=', 1)[1]
                        parsed_link_params = parse_qs(params_str)
                        if 'uddg' in parsed_link_params and parsed_link_params['uddg']:
                            link = unquote(parsed_link_params['uddg'][0])
                    except Exception as parse_e:
                        logging.warning(f"Erro ao parsear link do DuckDuckGo '{raw_link}': {parse_e}")

                resultados.append({
                    'titulo': title,
                    'snippet': snippet,
                    'link': link
                })

        if not resultados:
            return f"Nenhum resultado encontrado na web para '{query}'."

        return formatar_resposta({
            'query': query,
            'resultados': resultados,
            'total': len(resultados)
        })

    except requests.exceptions.Timeout:
        return f"Erro ao pesquisar na web: Timeout após {REQUEST_TIMEOUT} segundos."
    except Exception as e:
        logging.exception(f"Erro inesperado ao processar pesquisa web para '{query}'")
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def listar_arquivos(diretorio: str) -> str:
    """Lista os arquivos e diretórios contidos em um dado diretorio."""
    if not validar_caminho(diretorio):
        raise SecurityError(f"Acesso negado ao diretório: {diretorio}")
        
    try:
        conteudo = os.listdir(diretorio)
        if not conteudo:
            return f"O diretório '{diretorio}' está vazio."
            
        # Separar arquivos e diretórios
        arquivos = []
        diretorios = []
        
        for item in conteudo:
            caminho_completo = os.path.join(diretorio, item)
            if os.path.isfile(caminho_completo):
                arquivos.append(item)
            else:
                diretorios.append(item)
                
        return formatar_resposta({
            'diretorio': diretorio,
            'arquivos': sorted(arquivos),
            'diretorios': sorted(diretorios),
            'total_arquivos': len(arquivos),
            'total_diretorios': len(diretorios)
        })
        
    except Exception as e:
        logging.exception(f"Erro ao listar o diretório '{diretorio}'")
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def ler_arquivo(caminho_arquivo: str) -> str:
    """Lê o conteúdo de um arquivo de texto especificado pelo caminho."""
    if not validar_caminho(caminho_arquivo):
        raise SecurityError(f"Acesso negado ao arquivo: {caminho_arquivo}")
        
    if not validar_extensao(caminho_arquivo):
        raise SecurityError(f"Extensão de arquivo não permitida: {caminho_arquivo}")
        
    if not validar_tamanho_arquivo(caminho_arquivo):
        raise SecurityError(f"Arquivo muito grande: {caminho_arquivo}")
        
    try:
        with open(caminho_arquivo, 'r', encoding='utf-8') as f:
            conteudo = f.read()
            
        return formatar_resposta({
            'arquivo': caminho_arquivo,
            'tamanho': len(conteudo),
            'conteudo': conteudo
        })
        
    except Exception as e:
        logging.exception(f"Erro ao ler o arquivo '{caminho_arquivo}'")
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def escrever_arquivo(caminho_arquivo: str, conteudo: str) -> str:
    """Escreve o conteudo especificado em um arquivo. Sobrescreve se existir."""
    if not validar_caminho(caminho_arquivo):
        raise SecurityError(f"Acesso negado ao arquivo: {caminho_arquivo}")
        
    if not validar_extensao(caminho_arquivo):
        raise SecurityError(f"Extensão de arquivo não permitida: {caminho_arquivo}")
        
    if len(conteudo.encode('utf-8')) > MAX_FILE_SIZE:
        raise SecurityError(f"Conteúdo muito grande para escrever no arquivo")
        
    try:
        diretorio = os.path.dirname(caminho_arquivo)
        if diretorio and not os.path.exists(diretorio):
            os.makedirs(diretorio, exist_ok=True)

        with open(caminho_arquivo, 'w', encoding='utf-8') as f:
            f.write(conteudo)
            
        return formatar_resposta({
            'arquivo': caminho_arquivo,
            'tamanho': len(conteudo),
            'status': 'escrito com sucesso'
        })
        
    except Exception as e:
        logging.exception(f"Erro ao escrever no arquivo '{caminho_arquivo}'")
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def executar_comando(comando: str) -> str:
    """EXECUTA um comando shell. Retorna stdout e stderr. EXTREMA CAUTELA!"""
    if not any(cmd in comando for cmd in ALLOWED_COMMANDS):
        raise SecurityError(f"Comando não permitido: {comando}")
        
    try:
        result = subprocess.run(
            comando,
            shell=True,
            check=True,
            capture_output=True,
            text=True,
            encoding='utf-8',
            errors='replace',
            timeout=REQUEST_TIMEOUT
        )
        
        return formatar_resposta({
            'comando': comando,
            'stdout': result.stdout.strip(),
            'stderr': result.stderr.strip(),
            'returncode': result.returncode
        })
        
    except subprocess.TimeoutExpired:
        return f"Erro: Comando excedeu o timeout de {REQUEST_TIMEOUT} segundos"
    except Exception as e:
        logging.exception(f"Erro ao executar comando '{comando}'")
        return formatar_resposta({
            "erro": str(e),
            "tipo": type(e).__name__
        })

@tool
def calcular(expressao: str) -> str:
    """
    Calcula o resultado de uma expressão matemática simples.
    Suporta operações básicas e funções matemáticas comuns.
    """
    try:
        # Dicionário seguro de funções matemáticas
        safe_dict = {
            'sqrt': math.sqrt,
            'pow': math.pow,
            'sin': math.sin,
            'cos': math.cos,
            'tan': math.tan,
            'radians': math.radians,
            'degrees': math.degrees,
            'log': math.log,
            'log10': math.log10,
            'exp': math.exp,
            'pi': math.pi,
            'e': math.e
        }
        
        resultado = eval(expressao, {"__builtins__": None}, safe_dict)
        
        return formatar_resposta({
            'expressao': expressao,
            'resultado': resultado,
            'tipo': type(resultado).__name__
        })
        
    except Exception as e:
        logging.exception(f"Erro ao calcular expressão '{expressao}'")
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def criar_word(caminho_arquivo: str, texto: str) -> str:
    """Cria um arquivo Word (.docx) com o texto fornecido."""
    try:
        doc = Document()
        doc.add_paragraph(texto)
        doc.save(caminho_arquivo)
        return formatar_resposta({"arquivo": caminho_arquivo, "status": "criado com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def ler_word(caminho_arquivo: str) -> str:
    """Lê o texto de um arquivo Word (.docx)."""
    try:
        doc = Document(caminho_arquivo)
        texto = "\n".join([p.text for p in doc.paragraphs])
        return formatar_resposta({"arquivo": caminho_arquivo, "conteudo": texto})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def criar_excel(caminho_arquivo: str, dados: list) -> str:
    """Cria um arquivo Excel (.xlsx) com os dados fornecidos (lista de listas)."""
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        for linha in dados:
            ws.append(linha)
        wb.save(caminho_arquivo)
        return formatar_resposta({"arquivo": caminho_arquivo, "status": "criado com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def ler_excel(caminho_arquivo: str) -> str:
    """Lê o conteúdo de um arquivo Excel (.xlsx)."""
    try:
        wb = openpyxl.load_workbook(caminho_arquivo)
        ws = wb.active
        dados = [[cell.value for cell in row] for row in ws.iter_rows()]
        return formatar_resposta({"arquivo": caminho_arquivo, "conteudo": dados})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def criar_ppt(caminho_arquivo: str, titulos: list, textos: list) -> str:
    """Cria um arquivo PowerPoint (.pptx) com slides de títulos e textos."""
    try:
        prs = Presentation()
        for titulo, texto in zip(titulos, textos):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = titulo
            slide.placeholders[1].text = texto
        prs.save(caminho_arquivo)
        return formatar_resposta({"arquivo": caminho_arquivo, "status": "criado com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def ler_ppt(caminho_arquivo: str) -> str:
    """Lê os títulos e textos dos slides de um arquivo PowerPoint (.pptx)."""
    try:
        prs = Presentation(caminho_arquivo)
        slides = []
        for slide in prs.slides:
            titulo = slide.shapes.title.text if slide.shapes.title else ""
            texto = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    texto += shape.text + "\n"
            slides.append({"titulo": titulo, "texto": texto.strip()})
        return formatar_resposta({"arquivo": caminho_arquivo, "slides": slides})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def criar_arquivo_codigo(caminho_arquivo: str, codigo: str) -> str:
    """Cria um arquivo de código (ex: .py, .js, .html, etc) com o conteúdo fornecido."""
    try:
        diretorio = os.path.dirname(caminho_arquivo)
        if diretorio and not os.path.exists(diretorio):
            os.makedirs(diretorio, exist_ok=True)
        with open(caminho_arquivo, 'w', encoding='utf-8') as f:
            f.write(codigo)
        return formatar_resposta({"arquivo": caminho_arquivo, "status": "criado com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def criar_estrutura_pastas(base_dir: str, estrutura: dict) -> str:
    """
    Cria uma estrutura de pastas e arquivos a partir de um dicionário.
    Exemplo de estrutura:
    {
        "src": {
            "main.py": "print('Hello')",
            "utils": {
                "helpers.py": "# helpers"
            }
        },
        "README.md": "# Projeto"
    }
    """
    try:
        def criar_recursivo(base, tree):
            for nome, conteudo in tree.items():
                caminho = os.path.join(base, nome)
                if isinstance(conteudo, dict):
                    os.makedirs(caminho, exist_ok=True)
                    criar_recursivo(caminho, conteudo)
                else:
                    with open(caminho, 'w', encoding='utf-8') as f:
                        f.write(conteudo)
        criar_recursivo(base_dir, estrutura)
        return formatar_resposta({"base_dir": base_dir, "status": "estrutura criada com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def enviar_email(destinatario: str, assunto: str, mensagem: str) -> str:
    """Envia um e-mail para o destinatário especificado."""
    try:
        # Aqui você pode integrar com um serviço de e-mail como SMTP ou uma API
        # Exemplo usando SMTP:
        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart

        # Configurações do servidor SMTP
        smtp_server = "smtp.example.com"
        smtp_port = 587
        smtp_user = "seu_email@example.com"
        smtp_password = "sua_senha"

        # Criar mensagem
        msg = MIMEMultipart()
        msg['From'] = smtp_user
        msg['To'] = destinatario
        msg['Subject'] = assunto
        msg.attach(MIMEText(mensagem, 'plain'))

        # Enviar e-mail
        with smtplib.SMTP(smtp_server, smtp_port) as server:
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.send_message(msg)

        return formatar_resposta({"status": "E-mail enviado com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

@tool
def enviar_whatsapp(numero: str, mensagem: str) -> str:
    """Envia uma mensagem via WhatsApp para o número especificado."""
    try:
        # Aqui você pode integrar com uma API de WhatsApp como Twilio ou outra
        # Exemplo usando Twilio:
        from twilio.rest import Client

        # Configurações do Twilio
        account_sid = "AC7d476d8e5aad501b437f2de8c3c4bb13d"
        auth_token = "AC7d476d8e5aad501b437f2de8c3c4bb13:944d02ad5b70d5f506fd9d43859184cc"
        client = Client(account_sid, auth_token)

        # Enviar mensagem
        message = client.messages.create(
            body=mensagem,
            from_="whatsapp:+14155238886",  # Número do Twilio
            to=f"whatsapp:{numero}"
        )

        return formatar_resposta({"status": "Mensagem enviada com sucesso"})
    except Exception as e:
        return formatar_resposta({"erro": str(e), "tipo": type(e).__name__})

# Lista de todas as ferramentas disponíveis
ferramentas = [
    pesquisar_web,
    listar_arquivos,
    ler_arquivo,
    escrever_arquivo,
    executar_comando,
    calcular,
    criar_word,
    ler_word,
    criar_excel,
    ler_excel,
    criar_ppt,
    ler_ppt,
    criar_arquivo_codigo,
    criar_estrutura_pastas,
    enviar_email,
    enviar_whatsapp
] 